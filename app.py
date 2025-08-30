# Serve PDF files from the tmp directory
from flask import send_from_directory, send_file, Flask, request, render_template, jsonify
import os
import tempfile
import uuid
import logging
import threading
import requests
from pathlib import Path
import shutil
import time
from werkzeug.utils import secure_filename
from src.llm_providers import GeminiProvider, OpenAIProvider
from src.ppt_analyzer import PowerPointAnalyzer
from src.slide_generator import SlideGenerator
from src.utils import validate_file, cleanup_temp_files
import convertapi

# Configure logging
logging.basicConfig(level=logging.DEBUG)
logger = logging.getLogger(__name__)

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB max file size
app.config['UPLOAD_FOLDER'] = tempfile.gettempdir()

# Allowed file extensions
ALLOWED_EXTENSIONS = {'pptx', 'potx'}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

# Set your API key at the top of the file
convertapi.api_credentials = os.getenv("CLOUDCONVERT_API_KEY", "")

def convert_pdf_bg(session_dir):
    pptx_path = os.path.join(session_dir, 'output.pptx')
    pdf_path = os.path.join(session_dir, 'output.pdf')

    if not os.path.exists(pptx_path):
        logger.error(f"PPTX not found: {pptx_path}")
        return

    try:
        logger.info("Starting ConvertAPI PPTX → PDF conversion...")

        # Run conversion
        result = convertapi.convert(
            'pdf',
            {'File': pptx_path},
            from_format='pptx'
        )

        # Save converted PDF
        result.file.save(pdf_path)
        logger.info(f"PDF saved: {pdf_path}")

    except Exception as e:
        logger.error(f"Exception during ConvertAPI conversion: {e}")


# Run conversion in background
def start_pdf_conversion(session_dir):
    threading.Thread(target=convert_pdf_bg, args=(session_dir,), daemon=True).start()

@app.route('/')
def index():
    """Render the main page and clean up old tmp session directories"""

    tmp_root = Path(tempfile.gettempdir())
    now = time.time()
    max_age = 120  # seconds

    for session_dir in tmp_root.iterdir():
        if session_dir.is_dir() and session_dir.name.startswith("pptgen-"):
            try:
                mtime = session_dir.stat().st_mtime
                if now - mtime > max_age:
                    shutil.rmtree(session_dir)
                    logger.info(f"Deleted old tmp dir: {session_dir}")
            except Exception as e:
                logger.warning(f"Failed to delete old tmp dir {session_dir}: {e}")

    return render_template("index.html")

@app.route('/api/generate', methods=['POST'])
def generate_presentation():
    """Main endpoint to generate PowerPoint presentation"""
    try:
        # Validate request
        if 'template' not in request.files:
            return jsonify({'error': 'No template file uploaded'}), 400
        
        file = request.files['template']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        if not allowed_file(file.filename):
            return jsonify({'error': 'Invalid file type. Only .pptx and .potx files are allowed'}), 400

        # Get form data
        text_content = request.form.get('text_content', '').strip()
        guidance = request.form.get('guidance', '').strip()
        api_key = request.form.get('api_key', '').strip()
        ai_provider = request.form.get('ai_provider', 'gemini').lower().strip()
        ai_model = request.form.get('ai_model', '').strip()
        num_slides = request.form.get('num_slides', '').strip()
        # Always reuse images from template (made mandatory)
        reuse_images = True

        # --- TEXT CHARACTER LIMITATION ---
        MAX_TEXT_CHARS = 60000
        if not text_content:
            return jsonify({'error': 'Text content is required'}), 400
        if len(text_content) > MAX_TEXT_CHARS:
            text_content = text_content[:MAX_TEXT_CHARS]

        if not api_key:
            return jsonify({'error': 'API key is required'}), 400

        if ai_provider not in ['gemini', 'openai', 'aipipe']:
            return jsonify({'error': 'Invalid AI provider. Choose gemini, openai, or aipipe'}), 400

        # --- NUMBER OF SLIDES FUNCTIONALITY ---
        MIN_SLIDES = 7
        MAX_SLIDES = 40
        target_num_slides = None
        if num_slides:
            try:
                target_num_slides = max(MIN_SLIDES, min(MAX_SLIDES, int(num_slides)))
            except Exception:
                target_num_slides = MIN_SLIDES

                # Create unique tmp session dir
        session_id = f"pptgen-{uuid.uuid4()}"
        session_dir = os.path.join(tempfile.gettempdir(), session_id)
        os.makedirs(session_dir, exist_ok=True)

        # Save uploaded template
        filename = secure_filename(file.filename)
        template_path = os.path.join(session_dir, f"template_{filename}")
        file.save(template_path)

        # Validate the PowerPoint file
        if not validate_file(template_path):
            os.remove(template_path)
            return jsonify({'error': 'Invalid PowerPoint file'}), 400

        # Initialize components
        analyzer = PowerPointAnalyzer()
        generator = SlideGenerator()

        # Initialize the appropriate LLM provider
        if ai_provider == 'gemini':
            model_name = ai_model if ai_model else 'gemini-2.5-pro'
            llm_provider = GeminiProvider(api_key, model_name)
        elif ai_provider == 'openai':
            model_name = ai_model if ai_model else 'gpt-4o-mini'
            try:
                llm_provider = OpenAIProvider(api_key, model_name)
            except ImportError as e:
                return jsonify({'error': 'OpenAI library not available. Please install with: pip install openai'}), 400
        elif ai_provider == 'aipipe':
            model_name = ai_model if ai_model else 'openai/gpt-4o-mini'
            from src.llm_providers import AIPipeProvider
            llm_provider = AIPipeProvider(api_key, model_name)
        else:
            return jsonify({'error': f'Unsupported AI provider: {ai_provider}'}), 400

        # Process the presentation
        logger.info(f"Processing presentation for session {session_id} using {ai_provider}")

        # Step 1: Analyze template
        template_info = analyzer.analyze_template(template_path)

        # Define output path in session dir
        output_path = os.path.join(session_dir, 'output.pptx')

        # Step 2: Check if we need robust pipeline for multi-placeholder templates
        use_robust = False
        max_content_placeholders = 0

        # Analyze template for multiple content placeholders
        if template_info and 'existing_slides' in template_info:
            for slide in template_info.get('existing_slides', []):
                placeholders = slide.get('placeholders', [])
                content_phs = [p for p in placeholders 
                             if any(x in str(p.get('type', '')).upper() 
                                   for x in ['CONTENT', 'BODY', 'TEXT', 'OBJECT'])]
                if len(content_phs) > max_content_placeholders:
                    max_content_placeholders = len(content_phs)
                if len(content_phs) > 2:  # More than 2 content areas means multi-placeholder
                    use_robust = True

        # --- LLM GENERATION AND SLIDE COUNT ENFORCEMENT ---
        if use_robust:
            logger.info(f"Detected multi-placeholder template (max {max_content_placeholders} content areas), using robust pipeline")
            try:
                from src.robust_pipeline import RobustSlidePipeline
                pipeline = RobustSlidePipeline(llm_provider)
                refined_slides, selected_indices = pipeline.run(
                    text_content,
                    template_info,
                    template_info['presentation_object'],
                    output_path,
                    guidance,
                    num_slides=target_num_slides,
                    reuse_images=reuse_images
                )
                logger.info(f"Robust pipeline completed: {len(refined_slides)} slides generated")
            except Exception as e:
                logger.warning(f"Robust pipeline failed: {e}, falling back to standard generation")
                slide_structure = llm_provider.parse_text_to_slides(text_content, guidance, template_info, num_slides=target_num_slides)
                generator.create_presentation(slide_structure, template_info, output_path, use_robust_pipeline=False, reuse_images=reuse_images)
        else:
            logger.info("Using standard generation pipeline")
            slide_structure = llm_provider.parse_text_to_slides(text_content, guidance, template_info, num_slides=target_num_slides)
            generator.create_presentation(slide_structure, template_info, output_path, use_robust_pipeline=False, reuse_images=reuse_images)

        # Clean up template file
        os.remove(template_path)

        # --- Trigger PDF conversion in background ---
        start_pdf_conversion(session_dir)

        # Return success response with download info
        return jsonify({
            'success': True,
            'download_url': f'/api/download/{session_id}',
            'session_id': session_id
        })
        
    except Exception as e:
        logger.error(f"Error generating presentation: {str(e)}")
        # Clean up any temp files
        cleanup_temp_files(locals().get('template_path'), locals().get('output_path'))
        return jsonify({'error': f'Failed to generate presentation: {str(e)}'}), 500
    
# Endpoint to check if PDF is ready for preview
@app.route('/api/pdf_status/<session_id>')
def pdf_status(session_id):
    session_dir = os.path.join(tempfile.gettempdir(), session_id)
    pdf_path = os.path.join(session_dir, 'output.pdf')
    ready = os.path.exists(pdf_path)
    return jsonify({'ready': ready})

@app.route('/api/preview/<session_id>')
def preview_presentation(session_id):
    """Preview the generated presentation content"""
    try:
        file_path = os.path.join(tempfile.gettempdir(), session_id, 'output.pptx')
        
        if not os.path.exists(file_path):
            return jsonify({'error': 'File not found or expired'}), 404
        
        # Load the presentation to get slide info
        from pptx import Presentation
        pres = Presentation(file_path)
        
        slides_preview = []
        for i, slide in enumerate(pres.slides):
            slide_info = {
                'slide_number': i + 1,
                'title': '',
                'subtitle': '',
                'content': [],
                'layout_name': '',
                'slide_type': 'content'  # default
            }
            
            # Try to get layout name
            try:
                slide_info['layout_name'] = slide.slide_layout.name
            except:
                slide_info['layout_name'] = 'Unknown Layout'
            
            # Determine slide type based on layout or content
            layout_name = slide_info['layout_name'].lower()
            if 'title' in layout_name and ('only' in layout_name or i == 0):
                slide_info['slide_type'] = 'title'
            elif i == len(pres.slides) - 1 and ('conclusion' in layout_name or 'summary' in layout_name):
                slide_info['slide_type'] = 'conclusion'
            
            # Extract text from slide shapes, preserving hierarchy
            title_shapes = []
            content_shapes = []
            subtitle_shapes = []
            
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text = shape.text.strip()
                    
                    # Check if this is a placeholder and what type
                    if shape.is_placeholder:
                        try:
                            placeholder_type = str(shape.placeholder_format.type)
                            if 'TITLE' in placeholder_type.upper():
                                title_shapes.append(text)
                            elif 'SUBTITLE' in placeholder_type.upper():
                                subtitle_shapes.append(text)
                            elif 'CONTENT' in placeholder_type.upper() or 'BODY' in placeholder_type.upper():
                                content_shapes.append(text)
                            else:
                                content_shapes.append(text)
                        except:
                            # If we can't determine placeholder type, use heuristics
                            if not slide_info['title'] and len(text) < 100:
                                title_shapes.append(text)
                            else:
                                content_shapes.append(text)
                    else:
                        # Not a placeholder, use position and size heuristics
                        if not slide_info['title'] and len(text) < 100:
                            title_shapes.append(text)
                        else:
                            content_shapes.append(text)
            
            # Assign the extracted text
            if title_shapes:
                slide_info['title'] = title_shapes[0]  # Take first title
            if subtitle_shapes:
                slide_info['subtitle'] = subtitle_shapes[0]  # Take first subtitle
            
            # Process content shapes - split lines and clean up
            for content_text in content_shapes:
                if content_text != slide_info['title'] and content_text != slide_info['subtitle']:
                    # Split by lines and clean up
                    lines = [line.strip() for line in content_text.split('\n') if line.strip()]
                    slide_info['content'].extend(lines)
            
            # If no title found but we have content, use the first content line as title
            if not slide_info['title'] and slide_info['content']:
                slide_info['title'] = slide_info['content'].pop(0)
            
            slides_preview.append(slide_info)
        
        return jsonify({
            'success': True,
            'total_slides': len(slides_preview),
            'slides': slides_preview
        })
        
    except Exception as e:
        logger.error(f"Error previewing file: {str(e)}")
        return jsonify({'error': 'Failed to preview presentation'}), 500

@app.route('/api/download/<session_id>')
def download_presentation(session_id):
    file_path = os.path.join(tempfile.gettempdir(), session_id, 'output.pptx')
    if not os.path.exists(file_path):
        return jsonify({'error': 'File not found or expired'}), 404
    return send_file(file_path, as_attachment=True,
                     download_name='generated_presentation.pptx',
                     mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')

@app.route('/api/models')
def get_available_models():
    """Get available AI models for each provider"""
    try:
        models = {
            'gemini': {
                'models': [
                    'gemini-2.5-pro',
                    'gemini-1.5-pro',
                    'gemini-1.5-flash',
                    'gemini-1.0-pro',
                    'gemini-1.5'
                ],
                'default': 'gemini-2.5-pro'
            },
            'openai': {
                'models': OpenAIProvider.get_available_models(),
                'default': 'gpt-4o-mini'
            },
            'aipipe': {
                'models': [
                    'openai/gpt-4o-mini',
                    'openai/gpt-4o',
                    'anthropic/claude-3-5-sonnet',
                    'google/gemini-2.0-flash-exp',
                    'meta-llama/llama-3.1-70b-instruct'
                ],
                'default': 'openai/gpt-4o-mini'
            }
        }
        return jsonify(models)
    except Exception as e:
        logger.error(f"Error getting models: {str(e)}")
        return jsonify({'error': 'Failed to get available models'}), 500

@app.route('/api/health')
def health_check():
    """Health check endpoint"""
    return jsonify({'status': 'healthy'})

@app.errorhandler(413)
def too_large(e):
    return jsonify({'error': 'File too large. Maximum size is 50MB.'}), 413

@app.errorhandler(500)
def internal_error(error):
    logger.error(f"Internal server error: {str(error)}")
    return jsonify({'error': 'Internal server error'}), 500

# Create upload directory if it doesn't exist
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# Export app for Vercel
# Vercel will automatically detect this as the WSGI application

@app.route('/sessions/<session_id>/output.pdf')
def serve_pdf(session_id):
    session_dir = os.path.join(tempfile.gettempdir(), session_id)
    pdf_path = os.path.join(session_dir, 'output.pdf')
    if not os.path.exists(pdf_path):
        return 'PDF not found', 404
    return send_from_directory(session_dir, 'output.pdf', mimetype='application/pdf')

if __name__ == '__main__':
    # Run the app locally
    port = int(os.environ.get('PORT', 5000))
    app.run(debug=True, host='0.0.0.0', port=port)
