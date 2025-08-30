"""
Simple and effective slide content replacement
This module focuses on actually replacing placeholder text with generated content
"""

import logging
from typing import Dict, List, Any
from pptx.slide import Slide

try:
    from .multi_placeholder_handler import MultiPlaceholderHandler
except ImportError:
    try:
        from multi_placeholder_handler import MultiPlaceholderHandler
    except ImportError:
        MultiPlaceholderHandler = None

logger = logging.getLogger(__name__)

def replace_slide_content_simple(slide: Slide, content: Dict[str, Any]) -> bool:
    """
    Simple and direct slide content replacement with multi-placeholder support
    
    Args:
        slide: The slide to populate
        content: Dictionary with 'title', 'subtitle', and 'content' keys
        
    Returns:
        True if content was successfully placed, False otherwise
    """
    success = False
    
    try:
        # Log what we're trying to place
        logger.info(f"Replacing slide content: title='{content.get('title', 'No title')}'")
        
        # Check if we should use multi-placeholder handler
        if MultiPlaceholderHandler:
            content_list = content.get('content', [])
            
            # Check for multiple content placeholders or separator markers
            content_placeholders = MultiPlaceholderHandler.get_content_placeholders(slide)
            has_separators = any(
                MultiPlaceholderHandler.PLACEHOLDER_SEPARATOR in str(item) or 
                any(sep in str(item).upper() for sep in MultiPlaceholderHandler.ALTERNATIVE_SEPARATORS)
                for item in content_list if item
            )
            
            if len(content_placeholders) > 1 or has_separators:
                logger.info(f"Using multi-placeholder handler for slide with {len(content_placeholders)} content areas")
                return MultiPlaceholderHandler.replace_slide_content_multi_aware(slide, content)
        
        # Track which placeholders we've used
        used_placeholders = []
        
        # Step 1: Place the title
        title_text = content.get('title', '')
        if title_text and title_text.strip():
            for shape in slide.shapes:
                if shape.is_placeholder and shape not in used_placeholders:
                    try:
                        # Get placeholder type
                        ph_type = str(shape.placeholder_format.type)
                        
                        # Look for title placeholder
                        if 'TITLE' in ph_type.upper() or ph_type == '1':
                            # Clear and set text
                            shape.text = ""  # Clear first
                            shape.text = title_text[:100]  # Limit to 100 chars
                            used_placeholders.append(shape)
                            logger.debug(f"Placed title: {title_text[:50]}")
                            success = True
                            break
                    except Exception as e:
                        logger.debug(f"Could not use placeholder for title: {e}")
        
        # Step 2: Place the subtitle (if exists)
        subtitle_text = content.get('subtitle', '')
        if subtitle_text and subtitle_text.strip():
            for shape in slide.shapes:
                if shape.is_placeholder and shape not in used_placeholders:
                    try:
                        ph_type = str(shape.placeholder_format.type)
                        
                        # Look for subtitle or body placeholder
                        if 'SUBTITLE' in ph_type.upper() or 'SUB' in ph_type.upper() or ph_type == '2':
                            shape.text = ""  # Clear first
                            shape.text = subtitle_text[:150]  # Limit to 150 chars
                            used_placeholders.append(shape)
                            logger.debug(f"Placed subtitle: {subtitle_text[:50]}")
                            success = True
                            break
                    except Exception:
                        pass
        
        # Step 3: Place the main content (bullet points)
        content_items = content.get('content', [])
        if content_items and isinstance(content_items, list):
            # Filter valid content
            valid_items = [str(item).strip() for item in content_items if item and str(item).strip()]
            
            if valid_items:
                for shape in slide.shapes:
                    if shape.is_placeholder and shape not in used_placeholders:
                        try:
                            ph_type = str(shape.placeholder_format.type)
                            
                            # Look for content/body placeholder
                            if ('CONTENT' in ph_type.upper() or 
                                'BODY' in ph_type.upper() or 
                                'OBJECT' in ph_type.upper() or
                                ph_type in ['7', '14']):  # Common content placeholder types
                                
                                if shape.has_text_frame:
                                    # Clear the text frame
                                    shape.text_frame.clear()
                                    
                                    # Add each bullet point
                                    for i, item in enumerate(valid_items[:6]):  # Max 6 items
                                        if i == 0:
                                            # Use first paragraph
                                            p = shape.text_frame.paragraphs[0] if shape.text_frame.paragraphs else shape.text_frame.add_paragraph()
                                        else:
                                            # Add new paragraph
                                            p = shape.text_frame.add_paragraph()
                                        
                                        # Set the text (limit length to fit)
                                        p.text = item[:150]  # Limit each bullet to 150 chars
                                        p.level = 0  # Set as top-level bullet
                                    
                                    used_placeholders.append(shape)
                                    logger.debug(f"Placed {len(valid_items)} content items")
                                    success = True
                                    break
                        except Exception as e:
                            logger.debug(f"Could not use placeholder for content: {e}")
        
        # Step 4: Clear any remaining placeholders that weren't used
        for shape in slide.shapes:
            if shape.is_placeholder and shape not in used_placeholders:
                try:
                    if shape.has_text_frame:
                        # Check if it has placeholder text
                        if shape.text and ('click to' in shape.text.lower() or 
                                         'add' in shape.text.lower() or
                                         'text' in shape.text.lower() or
                                         'title' in shape.text.lower() or
                                         'subtitle' in shape.text.lower() or
                                         'content' in shape.text.lower()):
                            # Clear placeholder text
                            shape.text = ""
                            logger.debug(f"Cleared unused placeholder: {shape.text[:30] if shape.text else 'empty'}")
                except Exception:
                    pass
        
        return success
        
    except Exception as e:
        logger.error(f"Error in simple content replacement: {e}")
        return False


def clear_all_placeholder_text(slide: Slide):
    """
    Clear all placeholder text from a slide
    This ensures no template text remains
    BUT preserves numbered list markers (1, 2, 3, etc.)
    """
    try:
        for shape in slide.shapes:
            if shape.is_placeholder:
                try:
                    if shape.has_text_frame:
                        text = shape.text if shape.text else ""
                        text_lower = text.lower()
                        
                        # Check if this is a numbered list marker (just "1", "2", "3", etc.)
                        # These should be preserved as they indicate list structure
                        if text.strip() in ['1', '2', '3', '4', '5', '6', '7', '8', '9', '10', 
                                           '1.', '2.', '3.', '4.', '5.', '6.', '7.', '8.', '9.', '10.',
                                           'a', 'b', 'c', 'd', 'e', 'f',
                                           'a.', 'b.', 'c.', 'd.', 'e.', 'f.',
                                           'i', 'ii', 'iii', 'iv', 'v', 'vi','01', '02', '03', '04', '05', '06','07','08','09']:
                            logger.debug(f"Preserving numbered list marker: {text}")
                            continue  # Don't clear numbered list markers
                        
                        # Check if it looks like placeholder text
                        if any(keyword in text_lower for keyword in ['click', 'add', 'insert', 'type', 'placeholder', 'text here']):
                            shape.text = ""
                            logger.debug(f"Cleared placeholder text: {text[:50]}")
                except Exception:
                    pass
    except Exception as e:
        logger.error(f"Error clearing placeholder text: {e}")


def ensure_text_fits(shape, text: str, max_chars: int = 100) -> str:
    """
    Ensure text fits within the shape bounds
    
    Args:
        shape: The shape to fit text into
        text: The text to fit
        max_chars: Maximum characters allowed
        
    Returns:
        Truncated text that fits
    """
    if not text:
        return ""
    
    # Basic truncation
    if len(text) > max_chars:
        return text[:max_chars-3] + "..."
    
    return text
