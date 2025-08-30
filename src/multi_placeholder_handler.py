"""
Multi-placeholder content handler for PowerPoint slides.
Intelligently distributes content across multiple text placeholders on a single slide.
"""

import logging
from typing import Dict, List, Any, Optional, Tuple
from pptx.slide import Slide

logger = logging.getLogger(__name__)

class MultiPlaceholderHandler:
    """Handles content distribution across multiple text placeholders on a slide"""
    
    PLACEHOLDER_SEPARATOR = "[NEXT_PLACEHOLDER]"
    ALTERNATIVE_SEPARATORS = ["[PLACEHOLDER]", "---", "###", "[TEXT_AREA]"]
    
    @staticmethod
    def parse_multi_placeholder_content(content_list: List[str]) -> List[List[str]]:
        """
        Parse content list into multiple placeholder groups.
        
        Args:
            content_list: List of content items that may contain separator markers
            
        Returns:
            List of content groups, one for each placeholder
        """
        if not content_list:
            return [[]]
        
        # Define all possible separator patterns including numbered variations
        separator_patterns = [
            '[NEXT_PLACEHOLDER]', '[PLACEHOLDER]', '[PLACEHOLDER_', 
            '[TEXT_AREA]', '---', '###'
        ]
        
        # Check for separator markers
        placeholder_groups = [[]]
        current_group_index = 0
        
        for item in content_list:
            if not item:
                continue
            
            item_str = str(item)
            
            # Check if this entire item is a separator or contains separator pattern
            is_separator = False
            
            # First check for exact matches or patterns like [PLACEHOLDER_2], [PLACEHOLDER_3] etc.
            if any(sep in item_str.upper() for sep in ['[PLACEHOLDER', '[NEXT_PLACEHOLDER', '[TEXT_AREA']):
                is_separator = True
                # Start a new group
                placeholder_groups.append([])
                current_group_index += 1
                
                # Extract any text that's NOT part of the separator marker
                # Remove all variations of placeholder markers
                cleaned_text = item_str
                import re
                # Remove [PLACEHOLDER_X] patterns
                cleaned_text = re.sub(r'\[PLACEHOLDER[_\s]*\d*\]', '', cleaned_text, flags=re.IGNORECASE)
                cleaned_text = re.sub(r'\[NEXT_PLACEHOLDER\]', '', cleaned_text, flags=re.IGNORECASE)
                cleaned_text = re.sub(r'\[TEXT_AREA[_\s]*\d*\]', '', cleaned_text, flags=re.IGNORECASE)
                cleaned_text = cleaned_text.strip()
                
                # If there's remaining text after removing the marker, add it to the new group
                if cleaned_text:
                    placeholder_groups[current_group_index].append(cleaned_text)
            
            # Check for other separators like ---, ###
            elif item_str.strip() in ['---', '###']:
                is_separator = True
                placeholder_groups.append([])
                current_group_index += 1
            
            else:
                # This is regular content, add to current group
                placeholder_groups[current_group_index].append(item)
        
        # Remove empty groups
        placeholder_groups = [group for group in placeholder_groups if group]
        
        # If no separators were found, try to intelligently split content
        if len(placeholder_groups) == 1 and len(content_list) > 3:
            return MultiPlaceholderHandler._auto_split_content(content_list)
        
        return placeholder_groups if placeholder_groups else [[]]
    
    @staticmethod
    def _auto_split_content(content_list: List[str], max_placeholders: int = 4) -> List[List[str]]:
        """
        Automatically split content into groups for multiple placeholders.
        
        Args:
            content_list: List of content items
            max_placeholders: Maximum number of placeholder groups to create
            
        Returns:
            List of content groups
        """
        if not content_list:
            return [[]]
        
        # Clean content list
        clean_content = [item for item in content_list if item and str(item).strip()]
        if not clean_content:
            return [[]]
        
        num_items = len(clean_content)
        
        # IMPORTANT: Always split content when we have multiple placeholders
        # Even if we have few items, distribute them across placeholders
        
        if max_placeholders <= 1:
            return [clean_content]
        
        # Calculate items per placeholder
        items_per_placeholder = max(1, num_items // max_placeholders)
        remainder = num_items % max_placeholders
        
        groups = []
        start_idx = 0
        
        for i in range(max_placeholders):
            # Distribute items evenly, with extra items going to first groups
            items_for_this_group = items_per_placeholder + (1 if i < remainder else 0)
            
            if start_idx < num_items:
                end_idx = min(start_idx + items_for_this_group, num_items)
                group = clean_content[start_idx:end_idx]
                if group:  # Only add non-empty groups
                    groups.append(group)
                start_idx = end_idx
        
        # If we have fewer items than placeholders, create minimal groups
        if len(groups) < max_placeholders and num_items > 0:
            # Redistribute to ensure we use more placeholders
            if num_items >= max_placeholders:
                # We have enough items, distribute one per placeholder and extras to first
                groups = []
                for i in range(min(num_items, max_placeholders)):
                    groups.append([clean_content[i]])
                # Add remaining items to existing groups
                for i in range(max_placeholders, num_items):
                    groups[i % len(groups)].append(clean_content[i])
            else:
                # We have fewer items than placeholders, put one item per group
                groups = [[item] for item in clean_content]
        
        return groups if groups else [[]]
    
    @staticmethod
    def get_content_placeholders(slide: Slide) -> List[Any]:
        """
        Get all real content/body placeholders from a slide (excluding title, subtitle, and static markers like numbers/bullets).
        Args:
            slide: The slide to analyze
        Returns:
            List of content placeholder shapes (excluding static markers)
        """
        content_placeholders = []
        title_found = False
        subtitle_found = False
        static_marker_patterns = [
            r'^\d+$', r'^\d+\.$', r'^[ivxlc]+$', r'^[a-zA-Z]\.$', r'^•$', r'^-$', r'^–$', r'^\d{2}$'
        ]
        import re
        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue
            try:
                ph_type = str(shape.placeholder_format.type).upper()
                # Skip title placeholder (only the first one)
                if 'TITLE' in ph_type and not title_found:
                    title_found = True
                    continue
                # Skip subtitle placeholder (only the first one)
                if 'SUBTITLE' in ph_type and not subtitle_found:
                    subtitle_found = True
                    continue
                # Exclude static marker placeholders (numbers, bullets, etc.)
                if shape.has_text_frame:
                    text = shape.text.strip() if shape.text else ""
                    if any(re.fullmatch(pat, text) for pat in static_marker_patterns):
                        logger.debug(f"Skipping static marker placeholder: '{text}'")
                        continue
                # Collect content/body placeholders
                if shape.has_text_frame and ('CONTENT' in ph_type or 'BODY' in ph_type or 'TEXT' in ph_type):
                    content_placeholders.append(shape)
                elif shape.has_text_frame and 'OBJECT' in ph_type:
                    content_placeholders.append(shape)
            except Exception as e:
                logger.debug(f"Error checking placeholder type: {e}")
        logger.debug(f"Found {len(content_placeholders)} real content placeholders on slide")
        return content_placeholders
    
    @staticmethod
    def distribute_content_to_placeholders(slide: Slide, content_groups: List[List[str]]) -> bool:
        """
        Distribute content groups to available content placeholders.
        
        Args:
            slide: The slide to populate
            content_groups: List of content groups to distribute
            
        Returns:
            True if content was successfully distributed
        """
        try:
            # Get available content placeholders
            content_placeholders = MultiPlaceholderHandler.get_content_placeholders(slide)
            
            if not content_placeholders:
                logger.warning("No content placeholders found on slide")
                return False
            
            logger.info(f"Distributing {len(content_groups)} content groups to {len(content_placeholders)} placeholders")
            
            # Clear all content placeholders first
            for placeholder in content_placeholders:
                if placeholder.has_text_frame:
                    placeholder.text_frame.clear()
            
            # Distribute content groups to placeholders
            for i, (placeholder, content_group) in enumerate(zip(content_placeholders, content_groups)):
                if not content_group:
                    continue
                logger.debug(f"Filling placeholder {i+1} with {len(content_group)} items")
                # Populate the placeholder
                if placeholder.has_text_frame:
                    text_frame = placeholder.text_frame
                    # If the original text is a static marker, preserve it and only add content after
                    orig_text = placeholder.text.strip() if placeholder.text else ""
                    static_marker_patterns = [
                        r'^\d+$', r'^\d+\.$', r'^[ivxlc]+$', r'^[a-zA-Z]\.$', r'^•$', r'^-$', r'^–$', r'^\d{2}$'
                    ]
                    import re
                    is_static_marker = any(re.fullmatch(pat, orig_text) for pat in static_marker_patterns)
                    # Add content items
                    for j, item in enumerate(content_group):
                        if not item or str(item).strip() == "":
                            continue
                        item_text = str(item).strip()
                        
                        # Clean any remaining placeholder markers from the text
                        import re
                        item_text = re.sub(r'\[PLACEHOLDER[_\s]*\d*\]', '', item_text, flags=re.IGNORECASE)
                        item_text = re.sub(r'\[NEXT_PLACEHOLDER\]', '', item_text, flags=re.IGNORECASE)
                        item_text = re.sub(r'\[TEXT_AREA[_\s]*\d*\]', '', item_text, flags=re.IGNORECASE)
                        item_text = item_text.strip()
                        
                        # Skip if item became empty after cleaning or is a separator
                        if not item_text or item_text in ['---', '###']:
                            continue
                        if j == 0 and len(text_frame.paragraphs) > 0:
                            paragraph = text_frame.paragraphs[0]
                        else:
                            paragraph = text_frame.add_paragraph()
                        # If static marker, preserve it and append content
                        if is_static_marker and orig_text:
                            paragraph.text = f"{orig_text} {item_text}"
                        else:
                            paragraph.text = item_text
                        paragraph.level = 0  # Bullet level
            
            # Clear any remaining empty placeholders
            for i in range(len(content_groups), len(content_placeholders)):
                if i < len(content_placeholders):
                    placeholder = content_placeholders[i]
                    if placeholder.has_text_frame:
                        placeholder.text_frame.clear()
                        # Ensure at least one empty paragraph
                        if len(placeholder.text_frame.paragraphs) == 0:
                            placeholder.text_frame.add_paragraph()
            
            return True
            
        except Exception as e:
            logger.error(f"Error distributing content to placeholders: {e}")
            return False
    
    @staticmethod
    def replace_slide_content_multi_aware(slide: Slide, slide_data: Dict[str, Any]) -> bool:
        """
        Replace slide content with awareness of multiple text placeholders.
        
        Args:
            slide: The slide to update
            slide_data: The content data for the slide
            
        Returns:
            True if content was successfully replaced
        """
        try:
            # Handle title
            title_text = slide_data.get('title')
            if title_text and title_text.strip() and title_text.lower() != 'none':
                for shape in slide.shapes:
                    if shape.is_placeholder:
                        try:
                            ph_type = str(shape.placeholder_format.type).upper()
                            if 'TITLE' in ph_type:
                                shape.text = title_text
                                break
                        except Exception:
                            continue
            
            # Handle subtitle (for title slides)
            subtitle_text = slide_data.get('subtitle')
            if subtitle_text and subtitle_text.strip() and subtitle_text.lower() != 'none':
                for shape in slide.shapes:
                    if shape.is_placeholder:
                        try:
                            ph_type = str(shape.placeholder_format.type).upper()
                            if 'SUBTITLE' in ph_type:
                                shape.text = subtitle_text
                                break
                        except Exception:
                            continue
            
            # Handle content with multi-placeholder awareness
            content_list = slide_data.get('content', [])
            if content_list and isinstance(content_list, list):
                # Parse content into groups for multiple placeholders
                content_groups = MultiPlaceholderHandler.parse_multi_placeholder_content(content_list)
                
                # Check how many content placeholders we have
                content_placeholders = MultiPlaceholderHandler.get_content_placeholders(slide)
                
                # If we have multiple placeholders but only one content group, 
                # try to auto-split the content
                if len(content_placeholders) > 1 and len(content_groups) == 1:
                    logger.info(f"Auto-splitting content for {len(content_placeholders)} placeholders")
                    content_groups = MultiPlaceholderHandler._auto_split_content(
                        content_groups[0], 
                        max_placeholders=len(content_placeholders)
                    )
                
                # Distribute content to placeholders
                return MultiPlaceholderHandler.distribute_content_to_placeholders(slide, content_groups)
            
            return True
            
        except Exception as e:
            logger.error(f"Error in multi-aware content replacement: {e}")
            return False
