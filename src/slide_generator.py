import logging
from typing import Dict, List, Any, Optional
from pptx import Presentation
from pptx.slide import Slide
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io
import random
import re
try:
    from .content_mapper import ContentMapper
except ImportError:
    from content_mapper import ContentMapper

logger = logging.getLogger(__name__)

class SlideGenerator:
    """Generates PowerPoint slides from structured content and template information"""
    
    def __init__(self):
        self.presentation = None
        self.template_info = None

    def _fit_text_in_placeholder(self, placeholder, kind: str = 'content'):
        """Adjust font sizes slightly to reduce overflow. Conservative to preserve template styling."""
        try:
            if not placeholder.has_text_frame:
                return
            tf = placeholder.text_frame
            # Estimate size adjustments based on kind
            if kind == 'title':
                min_size, max_size = 24, 44
                step = -2
            elif kind == 'subtitle':
                min_size, max_size = 16, 28
                step = -1
            else:
                min_size, max_size = 14, 22
                step = -1
            
            # Try decreasing font size slightly if too many lines
            def count_lines():
                return sum(1 for p in tf.paragraphs if (p.text or '').strip())
            
            # Rough line limit based on shape height
            from pptx.util import Emu
            height_inches = placeholder.height / Emu(1 * 914400)
            rough_max_lines = max(3, int(height_inches * (1.8 if kind=='content' else 1.2)))

            # If lines exceed rough limit, shrink font a bit
            if count_lines() > rough_max_lines:
                for p in tf.paragraphs:
                    for run in p.runs:
                        if run.font.size and run.font.size.pt > min_size:
                            run.font.size = Pt(max(min_size, run.font.size.pt + step))
        except Exception as e:
            logger.debug(f"Text fit skipped due to: {e}")

    def _remove_empty_text_placeholders(self, slide: Slide, keep: Optional[List[Any]] = None):
        """Remove text placeholders that are empty. Optionally keep some placeholders.
        This avoids leaving stray empty text boxes on slides.
        """
        try:
            keep = [k for k in (keep or []) if k is not None]
            to_remove = []
            for shape in slide.shapes:
                try:
                    if not getattr(shape, 'is_placeholder', False):
                        continue
                    if keep and shape in keep:
                        continue
                    if getattr(shape, 'has_text_frame', False):
                        text_val = (shape.text or '').strip()
                        # Remove placeholder if it contains no text
                        if not text_val:
                            to_remove.append(shape)
                except Exception:
                    continue

            # Remove shapes collected
            for shp in to_remove:
                try:
                    slide.shapes._spTree.remove(shp._element)
                except Exception:
                    # If we cannot remove safely, just leave it
                    continue
        except Exception as e:
            logger.warning(f"Error removing empty placeholders: {e}")
    
    def create_presentation(self, slide_structure: List[Dict[str, Any]], 
                          template_info: Dict[str, Any], 
                          output_path: str,
                          use_robust_pipeline: bool = True,
                          llm_provider = None,
                          reuse_images: bool = False,
                          num_slides: int = None) -> None:
        """
        Create a new PowerPoint presentation by intelligently mapping content to template slides
        
        Args:
            slide_structure: List of slide data from LLM parsing
            template_info: Template analysis results
            output_path: Path where to save the generated presentation
        """
        try:
            logger.info(f"Creating presentation with {len(slide_structure)} slides")
            self.template_info = template_info
            self.presentation = template_info['presentation_object']

            # --- ENFORCE SLIDE COUNT ---
            if num_slides is not None:
                # If too few slides, pad with empty slides; if too many, truncate
                if len(slide_structure) < num_slides:
                    for i in range(num_slides - len(slide_structure)):
                        slide_structure.append({"slide_type": "content", "title": f"Extra Slide {len(slide_structure)+1}", "subtitle": "", "content": [""]})
                elif len(slide_structure) > num_slides:
                    slide_structure = slide_structure[:num_slides]

            # --- REUSE IMAGES FUNCTIONALITY ---
            template_pictures = []
            if reuse_images and self.presentation:
                for s in self.presentation.slides:
                    slide_specs = []
                    for shape in s.shapes:
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            try:
                                slide_specs.append({
                                    "blob": shape.image.blob,
                                    "left": int(shape.left), "top": int(shape.top),
                                    "width": int(shape.width), "height": int(shape.height),
                                })
                            except Exception:
                                pass
                    template_pictures.append(slide_specs)

            # Use ContentMapper to intelligently map content to template slides
            mapper = ContentMapper()
            mapped_content, selected_indices = mapper.map_content_to_template(
                slide_structure, template_info
            )

            existing_slides = list(self.presentation.slides)
            logger.info(f"Template has {len(existing_slides)} existing slides")
            logger.info(f"Selected {len(selected_indices)} slides for content")

            # Track which slides have been used
            used_slides = set(selected_indices)

            # Import the simple replacer
            try:
                from .simple_slide_replacer import replace_slide_content_simple, clear_all_placeholder_text
            except ImportError:
                from simple_slide_replacer import replace_slide_content_simple, clear_all_placeholder_text

            # Replace content in selected slides
            for idx, (content, slide_idx) in enumerate(zip(mapped_content, selected_indices)):
                if slide_idx < len(existing_slides):
                    existing_slide = existing_slides[slide_idx]
                    logger.info(f"Replacing content for slide {slide_idx+1}: {content.get('title', 'Untitled')}")
                    clear_all_placeholder_text(existing_slide)
                    success = replace_slide_content_simple(existing_slide, content)
                    if not success:
                        logger.warning(f"Failed to place content on slide {slide_idx+1}, trying fallback")
                        self._replace_slide_content(existing_slide, content)

                    # --- REUSE IMAGES: Copy images from template slide to generated slide ---
                    if reuse_images and idx < len(template_pictures):
                        for pic in template_pictures[idx]:
                            try:
                                # Insert image at same position/size as in template
                                existing_slide.shapes.add_picture(
                                    io.BytesIO(pic["blob"]),
                                    pic["left"], pic["top"],
                                    width=pic["width"], height=pic["height"]
                                )
                            except Exception:
                                pass

            # Remove unused slides (work backwards to avoid index issues)
            for i in range(len(existing_slides) - 1, -1, -1):
                if i not in used_slides:
                    logger.info(f"Removing unused template slide {i+1}")
                    self._remove_slide_by_index(i)

            # Save the presentation
            self.presentation.save(output_path)
            logger.info(f"Presentation saved to: {output_path} with {len(self.presentation.slides)} slides")
            
        except Exception as e:
            logger.error(f"Error creating presentation: {str(e)}")
            raise Exception(f"Failed to create presentation: {str(e)}")
    
    def _create_fresh_presentation_from_template(self, template_presentation: Presentation):
        """Create a fresh presentation by replacing template slide content"""
        try:
            logger.info("Using template as base and replacing slide content...")
            
            # Use the template directly to preserve all layouts and styling
            self.presentation = template_presentation
            
            # Store reference to existing slides so we can reuse them
            self.template_slides = list(self.presentation.slides)
            logger.info(f"Template has {len(self.template_slides)} existing slides")
            
            logger.info("Template setup completed successfully")
            
        except Exception as e:
            logger.error(f"Error setting up template: {e}")
            # Fallback: use template as-is
            self.presentation = template_presentation
            self.template_slides = []
    
    def _clear_all_slides_simple(self):
        """Clear all slides using a simple XML-based approach"""
        try:
            original_count = len(self.presentation.slides)
            logger.info(f"Clearing {original_count} slides from template...")
            
            # Clear the slide ID list directly
            slide_id_lst = self.presentation.slides._sldIdLst
            
            # Remove all slide IDs
            while len(slide_id_lst) > 0:
                slide_id_lst.remove(slide_id_lst[0])
            
            logger.info(f"Successfully cleared all slides. Now has {len(self.presentation.slides)} slides")
            
        except Exception as e:
            logger.error(f"Error clearing slides: {e}")
            # If clearing fails, continue anyway - new slides will be added
    
    def _copy_slide_master_and_layouts(self, template_presentation: Presentation):
        """Copy slide master and layouts from template to new presentation"""
        try:
            # Create a completely new presentation
            from pptx import Presentation
            import tempfile
            import os
            
            # Save the template first to ensure it's properly written
            temp_template_path = tempfile.mktemp(suffix='.pptx')
            template_presentation.save(temp_template_path)
            
            # Load it fresh as our working presentation
            self.presentation = Presentation(temp_template_path)
            
            # Now clear ALL existing slides properly
            original_slide_count = len(self.presentation.slides)
            logger.info(f"Template has {original_slide_count} slides, clearing them all...")
            
            # Remove slides using the proper method - work backwards
            slides = list(self.presentation.slides)
            for i in range(len(slides) - 1, -1, -1):
                try:
                    # Remove slide from the presentation's slide collection
                    slide_id = slides[i].slide_id
                    slide_rId = slides[i].part.partname.split('/')[-1].replace('.xml', '')
                    
                    # Remove from the slide ID list
                    sld_id_lst = self.presentation.slides._sldIdLst
                    for sld_id in list(sld_id_lst):
                        if sld_id.id == slide_id:
                            sld_id_lst.remove(sld_id)
                            break
                    
                    logger.debug(f"Removed slide {i+1}")
                except Exception as e:
                    logger.warning(f"Error removing slide {i}: {e}")
            
            # Clean up temp file
            try:
                os.unlink(temp_template_path)
            except:
                pass
            
            logger.info(f"Successfully cleared all slides, now has {len(self.presentation.slides)} slides")
            
        except Exception as e:
            logger.error(f"Error setting up presentation: {e}")
            # Ultimate fallback - use template as-is
            self.presentation = template_presentation
    
    def _clear_existing_slides(self):
        """Remove existing slides while preserving layouts and master slides"""
        try:
            # Method 1: Try to clear slides using python-pptx internal methods
            slides_to_remove = list(self.presentation.slides)
            
            # Remove slides in reverse order
            for slide in reversed(slides_to_remove):
                rId = self.presentation.part.relate_to(slide.part, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide")
                self.presentation.part.drop_rel(rId)
            
            # Clear the slide list
            slide_id_list = self.presentation.slides._sldIdLst
            slide_id_list.clear()
            
            logger.info(f"Cleared {len(slides_to_remove)} existing slides from template")
            
        except Exception as e:
            logger.warning(f"Method 1 failed, trying alternative method: {e}")
            try:
                # Method 2: Alternative approach using XML manipulation
                slide_id_list = self.presentation.slides._sldIdLst
                for _ in range(len(slide_id_list)):
                    slide_id_list.remove(slide_id_list[0])
                
                logger.info("Cleared existing slides using alternative method")
                
            except Exception as e2:
                logger.warning(f"Could not clear existing slides with any method: {e2}")
                # If we can't clear, we'll add new slides after existing ones
    
    def _create_slide(self, slide_data: Dict[str, Any]) -> Slide:
        """Create a single slide from slide data"""
        try:
            # Get the appropriate layout
            layout_index = self._get_layout_for_slide(slide_data['slide_type'])
            slide_layout = self.presentation.slide_layouts[layout_index]
            
            # Add slide
            slide = self.presentation.slides.add_slide(slide_layout)
            
            # Populate slide content
            self._populate_slide_content(slide, slide_data)
            
            # Apply template styling
            self._apply_template_styling(slide, slide_data)
            
            # Add images if appropriate
            self._add_template_images(slide, slide_data)
            
            logger.debug(f"Created slide: {slide_data['title']}")
            return slide
            
        except Exception as e:
            logger.error(f"Error creating slide '{slide_data.get('title', 'Unknown')}': {e}")
            raise
    
    def _get_layout_for_slide(self, slide_type: str) -> int:
        """Get the best layout index for a slide type"""
        layouts = self.template_info.get('slide_layouts', [])
        
        if not layouts:
            return 0
        
        # Mapping preferences for slide types
        layout_preferences = {
            'title': ['title', 'Title Slide', 'Title Only'],
            'content': ['content', 'Title and Content', 'Two Content', 'Content with Caption'],
            'conclusion': ['title', 'Title and Content', 'Title Only']
        }
        
        preferred_names = layout_preferences.get(slide_type, ['content'])
        
        # Find best matching layout by name
        for pref_name in preferred_names:
            for i, layout in enumerate(layouts):
                if pref_name.lower() in layout['name'].lower():
                    return i
        
        # Fallback: return first available layout or default
        return 0 if layouts else 0
    
    def _populate_slide_content(self, slide: Slide, slide_data: Dict[str, Any]):
        """Populate slide with text content"""
        try:
            slide_type = slide_data['slide_type']
            
            # Handle different slide types
            if slide_type == 'title':
                self._populate_title_slide(slide, slide_data)
            elif slide_type == 'content':
                self._populate_content_slide(slide, slide_data)
            elif slide_type == 'conclusion':
                self._populate_conclusion_slide(slide, slide_data)
            else:
                # Default to content slide
                self._populate_content_slide(slide, slide_data)
                
        except Exception as e:
            logger.warning(f"Error populating slide content: {e}")
            # Create minimal content within placeholders if content population fails
            self._safe_minimal_content(slide, slide_data)
    
    def _populate_title_slide(self, slide: Slide, slide_data: Dict[str, Any]):
        """Populate a title slide strictly using title/subtitle placeholders"""
        try:
            title_placeholder = None
            subtitle_placeholder = None
            used_placeholders = []

            # Identify title and subtitle placeholders explicitly
            for shape in slide.shapes:
                if shape.is_placeholder:
                    try:
                        ph_type = str(shape.placeholder_format.type).upper()
                        if 'TITLE' in ph_type:
                            # Prefer the first true TITLE placeholder encountered
                            if title_placeholder is None:
                                title_placeholder = shape
                        elif 'SUBTITLE' in ph_type:
                            if subtitle_placeholder is None:
                                subtitle_placeholder = shape
                    except Exception:
                        continue

            # Fallback: sometimes subtitle placeholder is labeled as CONTENT/BODY in certain templates
            if subtitle_placeholder is None:
                for shape in slide.shapes:
                    if shape.is_placeholder:
                        try:
                            ph_type = str(shape.placeholder_format.type).upper()
                            if ('CONTENT' in ph_type or 'BODY' in ph_type) and subtitle_placeholder is None:
                                subtitle_placeholder = shape
                        except Exception:
                            continue

            # Set title if available and not None
            title_text = slide_data.get('title')
            if title_placeholder and title_text and title_text.strip() and title_text.lower() != 'none':
                title_placeholder.text = title_text
                self._format_title_text(title_placeholder)
                used_placeholders.append(title_placeholder)

            # Set subtitle if available and not None
            subtitle_text = slide_data.get('subtitle')
            if subtitle_placeholder and subtitle_text and subtitle_text.strip() and subtitle_text.lower() != 'none':
                subtitle_placeholder.text = subtitle_text
                self._format_subtitle_text(subtitle_placeholder)
                used_placeholders.append(subtitle_placeholder)

            # Cleanup: remove any empty text placeholders on the slide
            self._remove_empty_text_placeholders(slide, keep=used_placeholders)

        except Exception as e:
            logger.warning(f"Error populating title slide: {e}")
    
    def _populate_content_slide(self, slide: Slide, slide_data: Dict[str, Any]):
        """Populate a content slide with title and bullet points, only in designated placeholders"""
        try:
            title_placeholder = None
            content_placeholder = None
            used_placeholders = []

            # Find placeholders
            for shape in slide.shapes:
                if shape.is_placeholder:
                    try:
                        placeholder_type = str(shape.placeholder_format.type).upper()
                        if 'TITLE' in placeholder_type:
                            if title_placeholder is None:
                                title_placeholder = shape
                        elif 'CONTENT' in placeholder_type or 'BODY' in placeholder_type:
                            if content_placeholder is None:
                                content_placeholder = shape
                    except Exception:
                        continue

            # Set title if not None
            title_text = slide_data.get('title')
            if title_placeholder and title_text and title_text.strip() and title_text.lower() != 'none':
                title_placeholder.text = title_text
                self._format_title_text(title_placeholder)
                used_placeholders.append(title_placeholder)

            # Set content if not None or empty
            content_list = slide_data.get('content')
            if content_placeholder and content_list and isinstance(content_list, list):
                # Filter out None or empty content items
                valid_content = [item for item in content_list 
                               if item and str(item).strip() and str(item).lower() != 'none']
                if valid_content:
                    self._populate_bullet_points(content_placeholder, valid_content)
                    used_placeholders.append(content_placeholder)

            # Cleanup: remove any empty text placeholders on the slide
            self._remove_empty_text_placeholders(slide, keep=used_placeholders)

        except Exception as e:
            logger.warning(f"Error populating content slide: {e}")
    
    def _populate_conclusion_slide(self, slide: Slide, slide_data: Dict[str, Any]):
        """Populate a conclusion slide"""
        # Conclusion slides are similar to content slides
        self._populate_content_slide(slide, slide_data)
    
    def _populate_bullet_points(self, placeholder, content_list: List[str]):
        """Add bullet points to a content placeholder preserving original formatting"""
        try:
            if not placeholder.has_text_frame:
                return
            
            text_frame = placeholder.text_frame
            
            # Store original formatting from first paragraph if exists
            original_format = None
            if len(text_frame.paragraphs) > 0:
                first_para = text_frame.paragraphs[0]
                original_format = {
                    'level': first_para.level,
                    'alignment': first_para.alignment,
                    'font_name': first_para.runs[0].font.name if first_para.runs else None,
                    'font_size': first_para.runs[0].font.size if first_para.runs else None,
                    'font_bold': first_para.runs[0].font.bold if first_para.runs else None
                }
            
            text_frame.clear()  # Clear existing content

            # Fit bullets based on estimated capacity if available
            max_lines = None
            chars_per_line = None
            try:
                # Try to derive capacity using template_info existing slide analysis if available
                # Not always possible here; do a rough estimate using shape dimensions
                from pptx.util import Emu
                width_inches = placeholder.width / Emu(1 * 914400)
                height_inches = placeholder.height / Emu(1 * 914400)
                chars_per_line = max(20, int(width_inches * 12))
                max_lines = max(3, int(height_inches * 2.5))
            except Exception:
                pass

            def chunk_text(text: str, max_len: int) -> List[str]:
                # Chunk text into readable segments without breaking words
                if max_len is None or len(text) <= max_len:
                    return [text]
                words = text.split()
                lines = []
                current = []
                cur_len = 0
                for w in words:
                    add = (1 if current else 0) + len(w)
                    if cur_len + add > max_len:
                        lines.append(' '.join(current))
                        current = [w]
                        cur_len = len(w)
                    else:
                        current.append(w)
                        cur_len += add
                if current:
                    lines.append(' '.join(current))
                return lines

            # Build paragraphs preserving format
            total_lines_used = 0
            for i, bullet in enumerate(content_list):
                if max_lines is not None and total_lines_used >= max_lines:
                    break
                    
                # Don't wrap if it's already formatted (numbered/bulleted)
                if re.match(r'^\d+[\.\)]\s', bullet):
                    # It's a numbered item - preserve exactly
                    lines_to_add = [bullet]
                else:
                    # Wrap long lines
                    lines_to_add = chunk_text(bullet.strip(), chars_per_line)
                
                for j, line in enumerate(lines_to_add):
                    if max_lines is not None and total_lines_used >= max_lines:
                        break
                        
                    if len(text_frame.paragraphs) == 0:
                        paragraph = text_frame.paragraphs[0]
                    else:
                        paragraph = text_frame.add_paragraph()
                    
                    paragraph.text = line
                    
                    # Apply original formatting if available
                    if original_format:
                        paragraph.level = original_format['level'] if j == 0 else original_format['level'] + 1
                        if original_format['alignment']:
                            paragraph.alignment = original_format['alignment']
                    else:
                        paragraph.level = 0 if j == 0 else 1
                    
                    # Apply text formatting
                    if paragraph.runs:
                        for run in paragraph.runs:
                            if original_format and original_format['font_name']:
                                run.font.name = original_format['font_name']
                            if original_format and original_format['font_size']:
                                run.font.size = original_format['font_size']
                    
                    total_lines_used += 1
                    
                if max_lines is not None and total_lines_used >= max_lines:
                    break

            # Finally, attempt to fit text size within the placeholder
            self._fit_text_in_placeholder(placeholder, kind='content')
        except Exception as e:
            logger.warning(f"Error adding bullet points: {e}")
    
    def _format_title_text(self, placeholder):
        """Apply minimal formatting to title text - preserve template formatting and fit text"""
        try:
            self._fit_text_in_placeholder(placeholder, kind='title')
        except Exception as e:
            logger.warning(f"Error formatting title text: {e}")
    
    def _format_subtitle_text(self, placeholder):
        """Apply formatting to subtitle text and fit text"""
        try:
            if not placeholder.has_text_frame:
                return
            for paragraph in placeholder.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.bold = False
                    default_font = self.template_info.get('fonts', {}).get('default_font', 'Calibri')
                    run.font.name = default_font
            # Fit the subtitle within bounds
            self._fit_text_in_placeholder(placeholder, kind='subtitle')
        except Exception as e:
            logger.warning(f"Error formatting subtitle text: {e})")
    
    def _format_bullet_text(self, paragraph):
        """Apply formatting to bullet point text"""
        try:
            for run in paragraph.runs:
                run.font.bold = False
                default_font = self.template_info.get('fonts', {}).get('default_font', 'Calibri')
                run.font.name = default_font
        except Exception as e:
            logger.warning(f"Error formatting bullet text: {e}")
    
    def _apply_template_styling(self, slide: Slide, slide_data: Dict[str, Any]):
        """Apply template colors and styling to the slide"""
        try:
            # This is where you would apply colors, fonts, and other styling
            # from the template analysis. For now, we rely on the layout's
            # built-in styling which is preserved from the template.
            
            sample_colors = self.template_info.get('theme_colors', {}).get('sample_colors', [])
            if sample_colors and random.random() > 0.7:  # Occasionally apply accent colors
                self._apply_accent_color(slide, sample_colors[0])
                
        except Exception as e:
            logger.warning(f"Error applying template styling: {e}")
    
    def _apply_accent_color(self, slide: Slide, color_hex: str):
        """Apply an accent color to text elements"""
        try:
            # Parse hex color
            if color_hex.startswith('#'):
                color_hex = color_hex[1:]
            
            if len(color_hex) == 6:
                r = int(color_hex[0:2], 16)
                g = int(color_hex[2:4], 16)  
                b = int(color_hex[4:6], 16)
                rgb_color = RGBColor(r, g, b)
                
                # Apply to some text elements
                for shape in slide.shapes:
                    if shape.has_text_frame and random.random() > 0.8:  # Occasionally
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.text.strip():  # Only if there's text
                                    run.font.color.rgb = rgb_color
                                    break  # Only first run
                            break  # Only first paragraph
                            
        except Exception as e:
            logger.warning(f"Error applying accent color: {e}")
    
    def _add_template_images(self, slide: Slide, slide_data: Dict[str, Any]):
        """Add images from template to appropriate slides"""
        try:
            template_images = self.template_info.get('images', [])
            
            # For title slides, try to add a representative image
            if slide_data['slide_type'] == 'title' and template_images:
                self._add_background_image(slide, template_images[0])
            
            # For content slides, occasionally add small decorative images
            elif slide_data['slide_type'] == 'content' and template_images and random.random() > 0.7:
                self._add_decorative_image(slide, random.choice(template_images))
                
        except Exception as e:
            logger.warning(f"Error adding template images: {e}")
    
    def _add_background_image(self, slide: Slide, image_info: Dict[str, Any]):
        """Add an image as a background or large element"""
        try:
            # Add image to slide
            image_stream = io.BytesIO(image_info['image_data'])
            
            # Position it in the background (bottom-right corner, smaller)
            left = Inches(8)  # Right side
            top = Inches(5)   # Bottom area
            width = Inches(2) # Smaller size
            height = Inches(1.5)
            
            slide.shapes.add_picture(image_stream, left, top, width, height)
            
        except Exception as e:
            logger.warning(f"Error adding background image: {e}")
    
    def _add_decorative_image(self, slide: Slide, image_info: Dict[str, Any]):
        """Add a small decorative image to the slide"""
        try:
            image_stream = io.BytesIO(image_info['image_data'])
            
            # Position it as a small decorative element
            left = Inches(8.5)  # Far right
            top = Inches(1)     # Top area
            width = Inches(1)   # Small size
            height = Inches(0.75)
            
            slide.shapes.add_picture(image_stream, left, top, width, height)
            
        except Exception as e:
            logger.warning(f"Error adding decorative image: {e}")
    
    def _replace_slide_content_preserving_format(self, slide: Slide, slide_data: Dict[str, Any]):
        """Replace content while preserving ALL template formatting and positioning"""
        try:
            logger.debug(f"Replacing content for slide with format preservation: {slide_data.get('title', 'Untitled')}")
            
            # Store original placeholder information before clearing
            placeholder_info = {}
            for shape in slide.shapes:
                if shape.is_placeholder:
                    try:
                        ph_type = str(shape.placeholder_format.type).upper()
                        if shape.has_text_frame:
                            # Store formatting from existing content
                            tf = shape.text_frame
                            if len(tf.paragraphs) > 0:
                                first_para = tf.paragraphs[0]
                                placeholder_info[ph_type] = {
                                    'shape': shape,
                                    'level': first_para.level,
                                    'alignment': first_para.alignment,
                                    'font_name': first_para.runs[0].font.name if first_para.runs else None,
                                    'font_size': first_para.runs[0].font.size if first_para.runs else None,
                                    'font_bold': first_para.runs[0].font.bold if first_para.runs else None,
                                    'original_text': shape.text
                                }
                    except Exception as e:
                        logger.debug(f"Could not store placeholder info: {e}")
            
            # Clear ALL placeholders completely before adding new content
            for shape in slide.shapes:
                if shape.is_placeholder and shape.has_text_frame:
                    # Clear the text frame completely
                    shape.text_frame.clear()
                    # Make sure we have at least one empty paragraph
                    if len(shape.text_frame.paragraphs) == 0:
                        shape.text_frame.add_paragraph()
            
            # Now populate with new content using stored formatting
            self._populate_with_preserved_format_impl(slide, slide_data, placeholder_info)
            
            logger.debug("Successfully replaced slide content with format preservation")
            
        except Exception as e:
            logger.error(f"Error replacing slide content with format preservation: {e}")
            # Fallback to regular replacement
            self._replace_slide_content(slide, slide_data)
    
    def _populate_with_preserved_format_impl(self, slide: Slide, slide_data: Dict[str, Any], 
                                      placeholder_info: Dict[str, Any]):
        """Implementation: Populate slide using preserved formatting info"""
        logger.debug(f"Populating slide with content: title='{slide_data.get('title')}', "
                    f"subtitle='{slide_data.get('subtitle')}', "
                    f"content items={len(slide_data.get('content', []))}")
        
        # First, clear ALL text from ALL placeholders to ensure no template text remains
        for shape in slide.shapes:
            if shape.is_placeholder and shape.has_text_frame:
                shape.text_frame.clear()
        
        # Handle title
        title_text = slide_data.get('title')
        if title_text and title_text.strip() and title_text.lower() != 'none':
            # Find title placeholder
            title_placed = False
            for shape in slide.shapes:
                if shape.is_placeholder:
                    try:
                        ph_type = str(shape.placeholder_format.type).upper()
                        if 'TITLE' in ph_type and not title_placed:
                            shape.text = title_text
                            # Apply formatting if we have it
                            if ph_type in placeholder_info:
                                info = placeholder_info[ph_type]
                                if shape.has_text_frame:
                                    for paragraph in shape.text_frame.paragraphs:
                                        if info.get('alignment'):
                                            paragraph.alignment = info['alignment']
                                        for run in paragraph.runs:
                                            if info.get('font_name'):
                                                run.font.name = info['font_name']
                                            if info.get('font_size'):
                                                run.font.size = info['font_size']
                            title_placed = True
                            break
                    except Exception:
                        continue
        
        # Handle subtitle
        subtitle_text = slide_data.get('subtitle')
        if subtitle_text and subtitle_text.strip() and subtitle_text.lower() != 'none':
            subtitle_placed = False
            for shape in slide.shapes:
                if shape.is_placeholder and not subtitle_placed:
                    try:
                        ph_type = str(shape.placeholder_format.type).upper()
                        # Check for subtitle or body placeholder (but not if already used for title)
                        if ('SUBTITLE' in ph_type or ('BODY' in ph_type and shape.text == '')):
                            shape.text = subtitle_text
                            # Apply formatting if we have it
                            if ph_type in placeholder_info:
                                info = placeholder_info[ph_type]
                                if shape.has_text_frame:
                                    for paragraph in shape.text_frame.paragraphs:
                                        if info.get('alignment'):
                                            paragraph.alignment = info['alignment']
                                        for run in paragraph.runs:
                                            if info.get('font_name'):
                                                run.font.name = info['font_name']
                                            if info.get('font_size'):
                                                run.font.size = info['font_size']
                            subtitle_placed = True
                            break
                    except Exception:
                        continue
        
        # Handle content
        content_list = slide_data.get('content')
        if content_list and isinstance(content_list, list):
            valid_content = [item for item in content_list 
                           if item and str(item).strip() and str(item).lower() != 'none']
            
            if valid_content:
                content_placed = False
                for shape in slide.shapes:
                    if shape.is_placeholder and not content_placed:
                        try:
                            ph_type = str(shape.placeholder_format.type).upper()
                            # Find content/body placeholder that hasn't been used
                            if ('CONTENT' in ph_type or 'BODY' in ph_type) and shape.text == '':
                                if shape.has_text_frame:
                                    text_frame = shape.text_frame
                                    text_frame.clear()
                                    
                                    # Add content with preserved formatting
                                    for i, item in enumerate(valid_content):
                                        if i == 0:
                                            p = text_frame.paragraphs[0]
                                        else:
                                            p = text_frame.add_paragraph()
                                        
                                        p.text = item
                                        
                                        # Apply formatting if we have it
                                        if ph_type in placeholder_info:
                                            info = placeholder_info[ph_type]
                                            p.level = info.get('level', 0)
                                            
                                            if info.get('alignment'):
                                                p.alignment = info['alignment']
                                            
                                            for run in p.runs:
                                                if info.get('font_name'):
                                                    run.font.name = info['font_name']
                                                if info.get('font_size'):
                                                    run.font.size = info['font_size']
                                
                                content_placed = True
                                break
                        except Exception:
                            continue
    
    def _replace_slide_content(self, slide: Slide, slide_data: Dict[str, Any]):
        """Replace content of an existing slide with new data, then clean up unused placeholders"""
        try:
            logger.debug(f"Replacing content for slide: {slide_data.get('title', 'Untitled')}")

            # Clear all existing text content from placeholders only (do not clear shapes that are not text)
            for shape in slide.shapes:
                try:
                    if shape.is_placeholder and shape.has_text_frame:
                        shape.text_frame.clear()
                except Exception:
                    continue

            # Populate with new content
            self._populate_slide_content(slide, slide_data)

            # Apply template styling
            self._apply_template_styling(slide, slide_data)

            # Remove any empty/unused text placeholders after population
            self._remove_empty_text_placeholders(slide)

            logger.debug("Successfully replaced slide content")

        except Exception as e:
            logger.error(f"Error replacing slide content: {e}")
            # If replacement fails, avoid creating extra text boxes; try a minimal safe content placement
            try:
                self._safe_minimal_content(slide, slide_data)
            except Exception as e2:
                logger.error(f"Minimal content placement also failed: {e2}")
    
    def _remove_slide_by_index(self, index: int):
        """Remove a slide by its index"""
        try:
            if index < 0 or index >= len(self.presentation.slides):
                logger.warning(f"Invalid slide index {index}, skipping removal")
                return
            
            # Get slide ID list
            slide_id_lst = self.presentation.slides._sldIdLst
            
            # Remove the slide at the specified index
            if index < len(slide_id_lst):
                slide_id_lst.remove(slide_id_lst[index])
                logger.debug(f"Removed slide at index {index}")
            
        except Exception as e:
            logger.error(f"Error removing slide at index {index}: {e}")
    
    def _safe_minimal_content(self, slide: Slide, slide_data: Dict[str, Any]):
        """Create basic text content if normal population fails.
        Note: We avoid adding ad-hoc text boxes to keep content within designated placeholders.
        """
        try:
            # Try to find any content/body placeholder and put minimal content there
            content_ph = None
            title_ph = None
            for shape in slide.shapes:
                if getattr(shape, 'is_placeholder', False):
                    try:
                        ph_type = str(shape.placeholder_format.type).upper()
                        if ('CONTENT' in ph_type or 'BODY' in ph_type) and content_ph is None:
                            content_ph = shape
                        if 'TITLE' in ph_type and title_ph is None:
                            title_ph = shape
                    except Exception:
                        continue

            # Minimal safe placement
            if title_ph and slide_data.get('title'):
                title_ph.text = slide_data['title']
            if content_ph:
                tf = content_ph.text_frame
                tf.clear()
                bullets = slide_data.get('content') or []
                if not bullets and slide_data.get('subtitle'):
                    bullets = [slide_data['subtitle']]
                if not bullets and slide_data.get('title'):
                    bullets = [slide_data['title']]
                if not bullets:
                    bullets = ["Content unavailable"]
                for i, item in enumerate(bullets[:5]):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = item
                    p.level = 0
            
        except Exception as e:
            logger.error(f"Error creating fallback content: {e}")
