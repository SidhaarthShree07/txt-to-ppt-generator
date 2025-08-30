import logging
from typing import Dict, List, Any, Tuple, Optional
from pptx import Presentation
from pptx.slide import Slide
from pptx.shapes.base import BaseShape
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
import tempfile
import os

logger = logging.getLogger(__name__)

class PowerPointAnalyzer:
    """Analyzes PowerPoint templates to extract styles, layouts, and assets"""
    
    def __init__(self):
        self.template_info = {}
    
    def analyze_template(self, template_path: str) -> Dict[str, Any]:
        """
        Analyze a PowerPoint template file and extract styling information
        
        Args:
            template_path: Path to the PowerPoint template file
            
        Returns:
            Dictionary containing template analysis results
        """
        try:
            logger.info(f"Analyzing template: {template_path}")
            
            # Load the presentation
            presentation = Presentation(template_path)
            
            analysis = {
                'slide_layouts': self._analyze_slide_layouts(presentation),
                'existing_slides': self._analyze_existing_slides(presentation),
                'theme_colors': self._extract_theme_colors(presentation),
                'fonts': self._extract_fonts(presentation),
                'images': self._extract_images(presentation),
                'slide_count': len(presentation.slides),
                'layout_count': len(presentation.slide_layouts),
                'master_slide': self._analyze_slide_master(presentation),
                'slide_dimensions': {
                    'width': presentation.slide_width,
                    'height': presentation.slide_height
                },
                'presentation_object': presentation  # Keep reference for generation
            }
            
            logger.info(f"Template analysis complete: {analysis['slide_count']} slides, {analysis['layout_count']} layouts")
            return analysis
            
        except Exception as e:
            logger.error(f"Error analyzing template: {str(e)}")
            raise Exception(f"Failed to analyze PowerPoint template: {str(e)}")
    
    def _analyze_slide_layouts(self, presentation: Presentation) -> List[Dict[str, Any]]:
        """Analyze available slide layouts"""
        layouts = []
        
        for i, layout in enumerate(presentation.slide_layouts):
            layout_info = {
                'index': i,
                'name': layout.name,
                'placeholders': self._analyze_placeholders(layout),
                'width': presentation.slide_width,
                'height': presentation.slide_height
            }
            layouts.append(layout_info)
        
        return layouts
    
    def _analyze_placeholders(self, layout) -> List[Dict[str, Any]]:
        """Analyze placeholders in a slide layout with detailed dimensions"""
        placeholders = []
        
        for placeholder in layout.placeholders:
            try:
                # Calculate text capacity based on placeholder dimensions
                from pptx.util import Emu, Pt
                
                width_inches = placeholder.width / Emu(1 * 914400)  # Convert to inches
                height_inches = placeholder.height / Emu(1 * 914400)
                
                # Estimate character capacity based on dimensions
                # Assuming average character width and line height
                chars_per_line = int(width_inches * 12)  # Rough estimate
                lines_capacity = int(height_inches * 3)  # Rough estimate
                
                placeholder_type = str(placeholder.placeholder_format.type)
                
                # Determine optimal text length based on placeholder type
                if 'TITLE' in placeholder_type.upper():
                    optimal_chars = min(60, chars_per_line)
                elif 'SUBTITLE' in placeholder_type.upper():
                    optimal_chars = min(100, chars_per_line * 2)
                else:  # Content/Body
                    optimal_chars = chars_per_line * lines_capacity
                
                placeholder_info = {
                    'index': placeholder.placeholder_format.idx,
                    'type': placeholder_type,
                    'name': getattr(placeholder, 'name', f'Placeholder {placeholder.placeholder_format.idx}'),
                    'left': placeholder.left,
                    'top': placeholder.top,
                    'width': placeholder.width,
                    'height': placeholder.height,
                    'width_inches': width_inches,
                    'height_inches': height_inches,
                    'chars_per_line': chars_per_line,
                    'lines_capacity': lines_capacity,
                    'optimal_char_count': optimal_chars
                }
                placeholders.append(placeholder_info)
            except Exception as e:
                logger.warning(f"Could not analyze placeholder: {e}")
                continue
        
        return placeholders
    
    def _extract_theme_colors(self, presentation: Presentation) -> Dict[str, Any]:
        """Extract theme colors from the presentation"""
        try:
            theme_part = presentation.part.package.part_related_by(
                "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
            )
            
            # Basic color extraction - this is simplified
            # In a full implementation, you'd parse the theme XML
            colors = {
                'has_theme': theme_part is not None,
                'extracted_colors': []
            }
            
            # Try to extract colors from existing slides
            if presentation.slides:
                sample_colors = self._sample_colors_from_slides(presentation)
                colors['sample_colors'] = sample_colors
            
            return colors
            
        except Exception as e:
            logger.warning(f"Could not extract theme colors: {e}")
            return {'has_theme': False, 'extracted_colors': []}
    
    def _sample_colors_from_slides(self, presentation: Presentation) -> List[str]:
        """Sample colors from existing slides"""
        colors = set()
        
        try:
            for slide in list(presentation.slides)[:3]:  # Sample first 3 slides
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, 'fill') and shape.fill.type is not None:
                            if hasattr(shape.fill, 'fore_color'):
                                color = shape.fill.fore_color.rgb
                                if color:
                                    colors.add(f"#{color}")
                    except Exception:
                        continue
        except Exception as e:
            logger.warning(f"Error sampling colors: {e}")
        
        return list(colors)[:10]  # Return up to 10 colors
    
    def _extract_fonts(self, presentation: Presentation) -> Dict[str, Any]:
        """Extract font information from the presentation"""
        fonts = set()
        
        try:
            # Sample fonts from existing slides
            for slide in list(presentation.slides)[:3]:  # Sample first 3 slides
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame'):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.font.name:
                                    fonts.add(run.font.name)
                                    
        except Exception as e:
            logger.warning(f"Error extracting fonts: {e}")
        
        return {
            'fonts_found': list(fonts),
            'default_font': list(fonts)[0] if fonts else 'Calibri'
        }
    
    def _extract_images(self, presentation: Presentation) -> List[Dict[str, Any]]:
        """Extract image information from the presentation"""
        images = []
        
        try:
            for slide_num, slide in enumerate(presentation.slides):
                slide_images = self._extract_images_from_slide(slide, slide_num)
                images.extend(slide_images)
                
        except Exception as e:
            logger.warning(f"Error extracting images: {e}")
        
        return images
    
    def _extract_images_from_slide(self, slide: Slide, slide_num: int) -> List[Dict[str, Any]]:
        """Extract images from a specific slide"""
        slide_images = []
        
        try:
            for shape_num, shape in enumerate(slide.shapes):
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_info = {
                            'slide_index': slide_num,
                            'shape_index': shape_num,
                            'left': shape.left,
                            'top': shape.top,
                            'width': shape.width,
                            'height': shape.height,
                            'image_data': shape.image.blob,
                            'filename': getattr(shape.image, 'filename', f'image_{slide_num}_{shape_num}.png')
                        }
                        slide_images.append(image_info)
                    except Exception as e:
                        logger.warning(f"Could not extract image from slide {slide_num}, shape {shape_num}: {e}")
                        continue
                        
        except Exception as e:
            logger.warning(f"Error processing slide {slide_num} for images: {e}")
        
        return slide_images
    
    def _analyze_slide_master(self, presentation: Presentation) -> Dict[str, Any]:
        """Analyze the slide master for default styling"""
        try:
            slide_master = presentation.slide_master
            
            master_info = {
                'width': presentation.slide_width,
                'height': presentation.slide_height,
                'background': self._analyze_background(slide_master),
                'placeholders': self._analyze_placeholders(slide_master)
            }
            
            return master_info
            
        except Exception as e:
            logger.warning(f"Could not analyze slide master: {e}")
            return {}
    
    def _analyze_background(self, slide_master) -> Dict[str, Any]:
        """Analyze background styling"""
        try:
            background_info = {
                'has_background': hasattr(slide_master, 'background'),
                'fill_type': None
            }
            
            if hasattr(slide_master, 'background') and slide_master.background:
                if hasattr(slide_master.background, 'fill'):
                    background_info['fill_type'] = str(slide_master.background.fill.type)
            
            return background_info
            
        except Exception as e:
            logger.warning(f"Error analyzing background: {e}")
            return {'has_background': False}
    
    def get_best_layout_for_slide_type(self, template_info: Dict[str, Any], slide_type: str) -> int:
        """
        Determine the best layout index for a given slide type
        
        Args:
            template_info: Template analysis results
            slide_type: Type of slide (title, content, conclusion)
            
        Returns:
            Index of the best matching layout
        """
        layouts = template_info.get('slide_layouts', [])
        
        if not layouts:
            return 0
        
        # Mapping preferences for slide types
        layout_preferences = {
            'title': ['title', 'Title Slide', 'Title Only'],
            'content': ['content', 'Title and Content', 'Two Content', 'Content with Caption'],
            'conclusion': ['title', 'Title and Content', 'Title Only']
        }
        
        preferred_names = layout_preferences.get(slide_type, ['content'])
        
        # Find best matching layout by name
        for pref_name in preferred_names:
            for i, layout in enumerate(layouts):
                if pref_name.lower() in layout['name'].lower():
                    return i
        
        # Fallback: return first available layout or default
        return 0 if layouts else 0
    
    def _analyze_existing_slides(self, presentation: Presentation) -> List[Dict[str, Any]]:
        """Analyze existing slides in the template with detailed placeholder information"""
        slides_info = []
        
        for slide_idx, slide in enumerate(presentation.slides):
            slide_info = {
                'slide_index': slide_idx,
                'layout_name': slide.slide_layout.name if hasattr(slide, 'slide_layout') else 'Unknown',
                'placeholders': [],
                'suggested_content_type': None,
                'has_title': False,
                'has_content': False,
                'has_subtitle': False,
                'content_format': None,  # Will detect: 'numbered_list', 'bullet_list', 'paragraph'
                'detected_patterns': []  # Store detected patterns for AI guidance
            }
            
            # Analyze each placeholder in the slide
            for shape in slide.shapes:
                if shape.is_placeholder:
                    try:
                        ph_type = str(shape.placeholder_format.type).upper()
                        
                        # Calculate dimensions for text fitting
                        from pptx.util import Emu
                        width_inches = shape.width / Emu(1 * 914400)
                        height_inches = shape.height / Emu(1 * 914400)
                        
                        # Estimate text capacity
                        if 'TITLE' in ph_type:
                            slide_info['has_title'] = True
                            max_chars = int(width_inches * 10)  # Titles are larger font
                            suggested_lines = 1
                        elif 'SUBTITLE' in ph_type:
                            slide_info['has_subtitle'] = True
                            max_chars = int(width_inches * 12)
                            suggested_lines = 2
                        elif 'CONTENT' in ph_type or 'BODY' in ph_type:
                            slide_info['has_content'] = True
                            max_chars = int(width_inches * 15)
                            suggested_lines = int(height_inches * 2.5)
                        else:
                            max_chars = int(width_inches * 12)
                            suggested_lines = int(height_inches * 2)
                        
                        # Analyze current text for format patterns
                        current_text = shape.text if hasattr(shape, 'text') else ''
                        format_info = self._analyze_text_format(shape, current_text)
                        
                        placeholder_data = {
                            'type': ph_type,
                            'index': shape.placeholder_format.idx,
                            'width_inches': width_inches,
                            'height_inches': height_inches,
                            'max_chars_per_line': max_chars,
                            'suggested_lines': suggested_lines,
                            'current_text': current_text,
                            'text_format': format_info['format'],
                            'list_style': format_info.get('list_style'),
                            'actual_text_length': format_info.get('text_length', 0),
                            'line_count': format_info.get('line_count', 0)
                        }
                        slide_info['placeholders'].append(placeholder_data)
                        
                        # Track detected patterns for content generation
                        if format_info['format'] and 'CONTENT' in ph_type:
                            slide_info['content_format'] = format_info['format']
                            if format_info.get('patterns'):
                                slide_info['detected_patterns'].extend(format_info['patterns'])
                    except Exception as e:
                        logger.warning(f"Error analyzing placeholder in slide {slide_idx}: {e}")
            
            # Determine suggested content type based on placeholders
            if slide_idx == 0 or (slide_info['has_title'] and slide_info['has_subtitle'] and not slide_info['has_content']):
                slide_info['suggested_content_type'] = 'title'
            elif slide_idx == len(presentation.slides) - 1 and slide_info['has_title']:
                slide_info['suggested_content_type'] = 'conclusion'
            elif slide_info['has_title'] and slide_info['has_content']:
                slide_info['suggested_content_type'] = 'content'
            else:
                slide_info['suggested_content_type'] = 'content'
            
            slides_info.append(slide_info)
        
        return slides_info
    
    def _analyze_text_format(self, shape, text: str) -> Dict[str, Any]:
        """Analyze text format to detect lists, paragraphs, and patterns"""
        format_info = {
            'format': None,
            'list_style': None,
            'patterns': [],
            'text_length': 0,
            'line_count': 0
        }
        
        if not text or not text.strip():
            return format_info
        
        format_info['text_length'] = len(text)
        
        # Check if shape has text frame with paragraphs
        if hasattr(shape, 'text_frame'):
            paragraphs = shape.text_frame.paragraphs
            format_info['line_count'] = len(paragraphs)
            
            # Analyze paragraph patterns
            has_bullets = False
            has_numbers = False
            avg_para_length = 0
            total_length = 0
            para_count = 0
            
            for para in paragraphs:
                para_text = para.text.strip()
                if not para_text:
                    continue
                    
                para_count += 1
                total_length += len(para_text)
                
                # Check for numbered list patterns
                import re
                if re.match(r'^\d+[.)\s]', para_text):
                    has_numbers = True
                    format_info['patterns'].append('numbered_item')
                elif re.match(r'^[a-z][.)\s]', para_text, re.IGNORECASE):
                    has_numbers = True
                    format_info['patterns'].append('lettered_item')
                elif re.match(r'^[•●▪▫◦‣⁃]', para_text):
                    has_bullets = True
                    format_info['patterns'].append('bullet_item')
                elif para.level > 0:  # Indented paragraph often means bullet
                    has_bullets = True
                    format_info['patterns'].append('indented_item')
            
            if para_count > 0:
                avg_para_length = total_length // para_count
            
            # Determine overall format
            if has_numbers:
                format_info['format'] = 'numbered_list'
                format_info['list_style'] = 'numbers'
            elif has_bullets or (para_count > 2 and avg_para_length < 100):
                format_info['format'] = 'bullet_list'
                format_info['list_style'] = 'bullets'
            elif para_count <= 2 and avg_para_length > 100:
                format_info['format'] = 'paragraph'
                format_info['list_style'] = None
            elif para_count > 0:
                format_info['format'] = 'mixed_content'
        else:
            # Simple text analysis without text frame
            lines = text.split('\n')
            format_info['line_count'] = len(lines)
            
            # Check for list patterns in plain text
            import re
            numbered_pattern = re.compile(r'^\d+[.)\s]')
            bullet_chars = ['•', '●', '▪', '▫', '◦', '‣', '⁃', '-', '*']
            
            has_numbers = any(numbered_pattern.match(line.strip()) for line in lines if line.strip())
            has_bullets = any(line.strip().startswith(tuple(bullet_chars)) for line in lines if line.strip())
            
            if has_numbers:
                format_info['format'] = 'numbered_list'
                format_info['list_style'] = 'numbers'
            elif has_bullets:
                format_info['format'] = 'bullet_list'
                format_info['list_style'] = 'bullets'
            elif len(lines) == 1 or (len(lines) <= 3 and all(len(l) > 50 for l in lines if l.strip())):
                format_info['format'] = 'paragraph'
            
        return format_info
